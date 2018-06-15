from pptx import Presentation
from glob import glob

# Presentationインスタンスの作成
ppt = Presentation()
# 幅
width = ppt.slide_width
print(width)
# 高さ
height = ppt.slide_height
print(height)

# レイアウト, 6番は白紙
blank_slide_layout = ppt.slide_layouts[5]

# 画像ファイルの読み込み
fnms = sorted(glob('./*point.png'))

# ファイル毎にループ

count = 0
for i, fnm in enumerate(fnms):
	# 白紙のスライドの追加
	if(i%4 == 0):
		slide = ppt.slides.add_slide(blank_slide_layout)
		count = 0

	# 画像の挿入
	pic = slide.shapes.add_picture(fnm, 0, 0)
	pic.width = int( width/2.0 )
	pic.height = int( height/2.0 )
	
	if(count == 0):
		pic.left = 0
		pic.top  = 0
	elif(count == 1):
		pic.left = width - pic.width
		pic.top  = 0
	elif(count == 2):
		pic.left = 0
		pic.top  = height - pic.height
	elif(count == 3):
		pic.left = width - pic.width
		pic.top  = height - pic.height
	else:
		print("error")

	count = count + 1

# 名前をつけて保存
ppt.save('figure.pptx')